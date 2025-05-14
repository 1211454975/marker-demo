from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from typing import List, Dict

class MoffeeAdapter:
    # 布局类型映射（对应PPT母版索引）
    layout_mapping = {
        "title-content": 1,  # 标题+内容布局
        "split": 3,          # 分栏布局
        "quote": 5           # 引用布局
    }
    
    def render_ppt(self, outline: List[Dict]):  # 修正缩进
        prs = Presentation()
        
        for slide_data in outline:
            slide_layout = self._select_layout(prs, slide_data["layout"])
            slide = prs.slides.add_slide(slide_layout)
            self._fill_content(slide, slide_data)
        
        return prs

    def _select_layout(self, prs, layout_type):
        """修复布局选择逻辑"""
        # 使用安全字典获取避免KeyError
        layout_idx = self.layout_mapping.get(
            layout_type.lower(),
            self.layout_mapping.get("title-content", 0)  # 双重fallback
        )
        
        # 添加母版索引有效性检查
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            raise ValueError(f"无效的布局索引: {layout_type} -> {layout_idx} (总布局数：{len(prs.slide_layouts)})")
            
        return layout

    def _fill_content(self, slide, data):
        title_shape = slide.shapes.title
        if title_shape and data.get("title"):
            title_shape.text = data["title"][:100]  # 标题长度限制

        # 统一内容处理逻辑
        content = data.get("content", "")
        placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx > 0]
        
        if placeholders:
            primary_ph = placeholders[0]
            primary_ph.text = content[:2000]  # 内容长度限制

        # 处理分栏布局
        if data["layout"] == "split":
            # 将内容分割为左右两栏
            parts = self._split_content(data["content"])
            placeholders = [
                ph for ph in slide.placeholders 
                if ph.placeholder_format.idx > 0  # 0是标题占位符
            ]
            
            for i, part in enumerate(parts[:2]):  # 最多填充两栏
                if i < len(placeholders):
                    placeholders[i].text = part

    def _split_content(self, content: str) -> list:
        """智能分割内容到分栏布局"""
        # 尝试按自然分隔符分割
        if "---" in content:
            return content.split("---", 1)
        
        # 自动平衡分割
        paragraphs = [p for p in content.split("\n") if p.strip()]
        mid = len(paragraphs) // 2
        return [
            "\n".join(paragraphs[:mid]),
            "\n".join(paragraphs[mid:])
        ]
